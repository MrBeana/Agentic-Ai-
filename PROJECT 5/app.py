import gradio as gr
from openai import OpenAI
from openai import OpenAIError
from dotenv import load_dotenv
import json
import os
import zipfile
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# ── OpenAI client ────────────────────────────────────────────────────────
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

MODEL = "gpt-5"

# ── Prompts ──────────────────────────────────────────────────────────────────
PROMPTS = {
    "api_docs": """You are an expert technical writer. Analyze the provided codebase and generate comprehensive API documentation.

Include:
- Overview of the API/module purpose
- All public functions/methods/classes with:
  - Description
  - Parameters (name, type, description, required/optional)
  - Return values (type, description)
  - Raises/Exceptions
  - Usage examples
- Authentication details (if applicable)
- Rate limits or constraints (if applicable)
- Code examples for common use cases

Format using Markdown with proper headers, tables, and code blocks.""",

    "readme": """You are an expert technical writer. Analyze the provided codebase and generate a professional README.md file.

Include:
- Project title and badge-ready header
- Concise project description (what it does, why it exists)
- Key features (bullet list)
- Tech stack / prerequisites
- Installation instructions (step-by-step)
- Quick start / usage example with code
- Configuration options (if any)
- Project structure overview
- Contributing guidelines
- License section

Make it welcoming, clear, and follow best practices from top open-source projects. Use Markdown formatting.""",

    "architecture": """You are a senior software architect. Analyze the provided codebase and generate a comprehensive architecture summary.

Include:
- High-level system overview
- Core components and their responsibilities
- Data flow diagrams (described in text/ASCII)
- Design patterns identified
- Module dependencies and interactions
- Technology choices and rationale (inferred)
- Scalability and maintainability observations
- Potential bottlenecks or areas for improvement
- Security considerations (if relevant)

Format using Markdown with clear sections, ASCII diagrams where helpful.""",

"ppt": """You are a senior software architect and presentation designer.

Analyze the provided codebase and generate slide content for a professional PowerPoint presentation.

Return the response in this EXACT JSON format:

{
  "title": "Project Title",
  "slides": [
    {
      "title": "Slide title",
      "content": [
        "bullet point 1",
        "bullet point 2"
      ]
    }
  ]
}

Rules:
- Generate 6-10 slides
- Keep bullets concise
- Focus on architecture, features, APIs, flow, and usage
- Output ONLY valid JSON
""",
}

# ── Core generation function ─────────────────────────────────────────────────
def generate_documentation(code_input: str, file_obj, doc_type: str, extra_context: str):
    """Stream documentation from OpenAIAPI."""
    
    # Resolve code source
    code_content = ""
    file_info = ""
    
    if file_obj is not None:
        try:
            with open(file_obj.name, "r", encoding="utf-8", errors="replace") as f:
                code_content = f.read()
            file_info = f"\n**Source file:** `{Path(file_obj.name).name}`\n"
        except Exception as e:
            code_content = f"[Error reading file: {e}]"
    elif code_input.strip():
        code_content = code_input.strip()
    else:
        yield "⚠️ Please paste code or upload a file first.", "", ""
        return

    if not code_content.strip():
        yield "⚠️ The provided file appears to be empty.", "", ""
        return

    system_prompt = PROMPTS.get(doc_type, PROMPTS["readme"])
    
    user_message = f"""Please analyze this codebase and generate the requested documentation.
{file_info}
{f'Additional context: {extra_context}' if extra_context.strip() else ''}

```
{code_content[:15000]}
```
{'[Note: Code was truncated to 15,000 characters]' if len(code_content) > 15000 else ''}
"""

    # Stream response
    output = ""
    try:
        stream = client.chat.completions.create(
            model=MODEL,
            messages=[
            {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_message},],
    max_completion_tokens=4096,
    stream=True,)

        for chunk in stream:
            if chunk.choices and chunk.choices[0].delta.content:
                text = chunk.choices[0].delta.content
                output += text
                yield output, "", ""
        # PPT handling
        if doc_type == "ppt":
            slide_json = json.loads(output)
            ppt_path = create_ppt(slide_json)

            yield (
                "✅ PowerPoint generated successfully",
                output,
                "PPT ready",
                output,
                ppt_path,
            )
            return
        # After streaming, build download content
        doc_label = {"api_docs": "API Documentation", "readme": "README", "architecture": "Architecture Summary"}[doc_type]
        yield output, output, f"✅ {doc_label} generated — {len(output):,} characters"
    
    except OpenAIError as e:
        yield f"❌ OpenAI API error: {e}", "", ""
    except Exception as e:
        yield f"❌ Unexpected error: {e}", "", ""


def create_download_file(markdown_content: str, doc_type: str):
    """Write markdown to a temp file and return path for Gradio download."""
    if not markdown_content.strip():
        return None
    
    filenames = {"api_docs": "API_DOCUMENTATION.md", "readme": "README.md", "architecture": "ARCHITECTURE.md"}
    fname = filenames.get(doc_type, "documentation.md")
    
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".md", mode="w", encoding="utf-8", prefix=fname.replace(".md", "_"))
    tmp.write(markdown_content)
    tmp.close()
    return tmp.name


def load_sample(sample_name: str):
    """Load a sample code file from sample_data/."""
    samples_dir = Path(__file__).parent / "sample_data"
    path = samples_dir / sample_name
    if path.exists():
        return path.read_text(encoding="utf-8")
    return ""

def create_ppt(slide_data: dict):
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = slide_data.get("title", "Project Presentation")
    subtitle.text = "Generated by DocGen AI"

    # Content Slides
    bullet_layout = prs.slide_layouts[1]

    for s in slide_data.get("slides", []):
        slide = prs.slides.add_slide(bullet_layout)

        title = slide.shapes.title
        title.text = s.get("title", "Slide")

        body = slide.placeholders[1]
        tf = body.text_frame

        bullets = s.get("content", [])

        if bullets:
            tf.text = bullets[0]

            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet

    tmp = tempfile.NamedTemporaryFile(
        delete=False,
        suffix=".pptx"
    )

    prs.save(tmp.name)

    return tmp.name


# ── CSS ──────────────────────────────────────────────────────────────────────
CSS = """
@import url('https://fonts.googleapis.com/css2?family=Space+Mono:wght@400;700&family=Syne:wght@400;600;800&display=swap');

:root {
  --bg: #0a0a0f;
  --surface: #111118;
  --border: #1e1e2e;
  --accent: #7c3aed;
  --accent2: #06b6d4;
  --accent3: #f59e0b;
  --text: #e2e8f0;
  --muted: #64748b;
  --success: #10b981;
  --mono: 'Space Mono', monospace;
  --sans: 'Syne', sans-serif;
}

* { box-sizing: border-box; }

body, .gradio-container {
  background: var(--bg) !important;
  font-family: var(--sans) !important;
  color: var(--text) !important;
}

/* ── Header ── */
.app-header {
  text-align: center;
  padding: 3rem 1rem 2rem;
  position: relative;
}
.app-header::before {
  content: '';
  position: absolute;
  top: 0; left: 50%; transform: translateX(-50%);
  width: 600px; height: 300px;
  background: radial-gradient(ellipse at center, rgba(124,58,237,.18) 0%, transparent 70%);
  pointer-events: none;
}
.app-title {
  font-size: 3rem;
  font-weight: 800;
  letter-spacing: -1px;
  background: linear-gradient(135deg, #7c3aed, #06b6d4, #f59e0b);
  -webkit-background-clip: text;
  -webkit-text-fill-color: transparent;
  background-clip: text;
  margin: 0;
  line-height: 1.1;
}
.app-subtitle {
  color: var(--muted);
  font-size: 1rem;
  margin-top: .5rem;
  font-family: var(--mono);
  letter-spacing: 2px;
  text-transform: uppercase;
}
.badge-row {
  display: flex;
  gap: .5rem;
  justify-content: center;
  margin-top: 1rem;
  flex-wrap: wrap;
}
.badge {
  background: var(--border);
  border: 1px solid var(--accent);
  color: var(--accent2);
  font-family: var(--mono);
  font-size: .7rem;
  padding: .25rem .75rem;
  border-radius: 999px;
  letter-spacing: 1px;
}

/* ── Panels ── */
.panel {
  background: var(--surface) !important;
  border: 1px solid var(--border) !important;
  border-radius: 12px !important;
  padding: 1.25rem !important;
}

/* ── Doc type cards ── */
.doc-cards {
  display: grid;
  grid-template-columns: repeat(3, 1fr);
  gap: .75rem;
  margin: .5rem 0 1.5rem;
}
.doc-card {
  background: var(--surface);
  border: 1px solid var(--border);
  border-radius: 10px;
  padding: 1rem;
  cursor: pointer;
  transition: all .2s ease;
  text-align: center;
}
.doc-card:hover, .doc-card.active {
  border-color: var(--accent);
  background: rgba(124,58,237,.1);
  transform: translateY(-2px);
}
.doc-card .icon { font-size: 1.5rem; }
.doc-card .label { font-size: .85rem; font-weight: 600; margin-top: .25rem; color: var(--text); }
.doc-card .desc { font-size: .7rem; color: var(--muted); margin-top: .2rem; }

/* ── Gradio overrides ── */
.gr-button-primary {
  background: linear-gradient(135deg, var(--accent), var(--accent2)) !important;
  border: none !important;
  color: #fff !important;
  font-family: var(--sans) !important;
  font-weight: 700 !important;
  letter-spacing: .5px !important;
  border-radius: 8px !important;
  transition: opacity .2s !important;
}
.gr-button-primary:hover { opacity: .85 !important; }

.gr-button-secondary {
  background: transparent !important;
  border: 1px solid var(--border) !important;
  color: var(--muted) !important;
  font-family: var(--mono) !important;
  font-size: .8rem !important;
  border-radius: 8px !important;
}
.gr-button-secondary:hover {
  border-color: var(--accent2) !important;
  color: var(--accent2) !important;
}

textarea, .gr-textbox textarea {
  background: #0d0d14 !important;
  border: 1px solid var(--border) !important;
  color: var(--text) !important;
  font-family: var(--mono) !important;
  font-size: .82rem !important;
  border-radius: 8px !important;
  line-height: 1.6 !important;
}
textarea:focus { border-color: var(--accent) !important; outline: none !important; }

.gr-radio label, .gr-checkbox label { color: var(--text) !important; font-family: var(--sans) !important; }

/* Tab styling */
.tab-nav button {
  background: transparent !important;
  border-bottom: 2px solid transparent !important;
  color: var(--muted) !important;
  font-family: var(--sans) !important;
  font-weight: 600 !important;
  border-radius: 0 !important;
  transition: all .2s !important;
}
.tab-nav button.selected {
  color: var(--accent2) !important;
  border-bottom-color: var(--accent2) !important;
  background: transparent !important;
}

/* Status bar */
.status-bar {
  font-family: var(--mono);
  font-size: .75rem;
  color: var(--success);
  background: rgba(16,185,129,.08);
  border: 1px solid rgba(16,185,129,.2);
  border-radius: 6px;
  padding: .4rem .8rem;
}

/* Section labels */
.section-label {
  font-family: var(--mono);
  font-size: .7rem;
  letter-spacing: 2px;
  text-transform: uppercase;
  color: var(--muted);
  margin-bottom: .5rem;
}

label span { color: var(--muted) !important; font-family: var(--mono) !important; font-size: .75rem !important; }

/* Output area */
.output-md {
  background: #0d0d14 !important;
  border-radius: 10px !important;
  padding: 1.25rem !important;
  min-height: 300px;
}
.output-md p, .output-md li { color: var(--text) !important; line-height: 1.7 !important; }
.output-md h1, .output-md h2, .output-md h3 { color: #fff !important; }
.output-md code { background: var(--border) !important; color: var(--accent2) !important; border-radius: 4px !important; padding: .1em .4em !important; }
.output-md pre { background: var(--border) !important; border-radius: 8px !important; padding: 1rem !important; }
.output-md a { color: var(--accent) !important; }

/* Upload area */
.gr-file-upload {
  border: 2px dashed var(--border) !important;
  border-radius: 10px !important;
  background: #0d0d14 !important;
}
.gr-file-upload:hover { border-color: var(--accent) !important; }

/* Footer */
.app-footer {
  text-align: center;
  padding: 2rem 1rem;
  color: var(--muted);
  font-family: var(--mono);
  font-size: .7rem;
  letter-spacing: 1px;
  border-top: 1px solid var(--border);
  margin-top: 2rem;
}

/* Scrollbar */
::-webkit-scrollbar { width: 6px; }
::-webkit-scrollbar-track { background: var(--bg); }
::-webkit-scrollbar-thumb { background: var(--border); border-radius: 3px; }
::-webkit-scrollbar-thumb:hover { background: var(--accent); }

/* Sample buttons row */
.sample-row { display: flex; gap: .5rem; flex-wrap: wrap; }

/* animate in */
@keyframes fadeUp {
  from { opacity: 0; transform: translateY(20px); }
  to   { opacity: 1; transform: translateY(0); }
}
.gradio-container > * { animation: fadeUp .4s ease forwards; }
"""

# ── UI ───────────────────────────────────────────────────────────────────────
def build_ui():
    with gr.Blocks(css=CSS, title="AI Documentation Generator", theme=gr.themes.Base()) as demo:

        # ── Header ──────────────────────────────────────────────────────────
        gr.HTML("""
        <div class="app-header">
          <h1 class="app-title">⚡ DocGen AI</h1>
          <p class="app-subtitle">Instant Documentation from Any Codebase</p>
          <div class="badge-row">
            <span class="badge">Claude Opus</span>
            <span class="badge">API Docs</span>
            <span class="badge">README</span>
            <span class="badge">Architecture</span>
            <span class="badge">Streaming</span>
          </div>
        </div>
        """)

        # ── Main layout ──────────────────────────────────────────────────────
        with gr.Row(equal_height=False):

            # ── Left: Input ──────────────────────────────────────────────────
            with gr.Column(scale=5):
                gr.HTML('<p class="section-label">📁 Code Input</p>')

                with gr.Tabs() as input_tabs:
                    with gr.TabItem("✏️  Paste Code"):
                        code_input = gr.Textbox(
                            placeholder="Paste your code here — Python, JS, Go, Rust, anything...",
                            lines=20,
                            max_lines=40,
                            label="",
                            show_copy_button=True,
                        )
                        with gr.Row():
                            sample_py  = gr.Button("🐍 Python API sample", size="sm", variant="secondary")
                            sample_js  = gr.Button("🟨 JS Class sample",   size="sm", variant="secondary")
                            sample_go  = gr.Button("🔵 Go Service sample", size="sm", variant="secondary")

                    with gr.TabItem("📂  Upload File"):
                        file_upload = gr.File(
                            label="Drop a source file (.py .js .ts .go .rs .java ...)",
                            file_types=[".py",".js",".ts",".go",".rs",".java",".rb",".cpp",".c",".cs",".php",".kt",".swift",".txt",".md"],
                        )

                gr.HTML('<p class="section-label" style="margin-top:1.5rem">📋 Documentation Type</p>')
                doc_type = gr.Radio(
                    choices=[("📄  API Docs", "api_docs"), ("📖  README.md", "readme"), ("🏛  Architecture", "architecture"),("📊 PowerPoint", "ppt")],
                    value="readme",
                    label="",
                    interactive=True,
                )

                extra_context = gr.Textbox(
                    placeholder="Optional: Add context (e.g. 'This is a FastAPI backend for a SaaS billing system')",
                    lines=2,
                    label="💬 Extra Context (optional)",
                )

                with gr.Row():
                    generate_btn = gr.Button("⚡ Generate Documentation", variant="primary", scale=4)
                    clear_btn    = gr.Button("🗑 Clear", variant="secondary", scale=1)

                status_out = gr.Textbox(label="", show_label=False, interactive=False, elem_classes=["status-bar"])

            # ── Right: Output ─────────────────────────────────────────────────
            with gr.Column(scale=7):
                gr.HTML('<p class="section-label">📝 Generated Documentation</p>')

                with gr.Tabs():
                    with gr.TabItem("👁  Preview"):
                        md_out = gr.Markdown(
                            value="*Your generated documentation will appear here...*",
                            elem_classes=["output-md"],
                        )
                    with gr.TabItem("🗒  Raw Markdown"):
                        raw_out = gr.Textbox(
                            value="",
                            lines=30,
                            label="",
                            show_copy_button=True,
                            interactive=False,
                        )

                with gr.Row():
                    download_btn = gr.Button("⬇️  Download .md", variant="secondary")
                    download_file = gr.File(label="Download", visible=True)
                    ppt_download = gr.File(label="Download PPT",visible=True)
        # ── Footer ───────────────────────────────────────────────────────────
        gr.HTML("""
        <div class="app-footer">
          DOCGEN AI &nbsp;|&nbsp; Powered by OpenAI gpt-5 &amp; Gradio &nbsp;|&nbsp; 
          Generate · API Docs · README · Architecture
        </div>
        """)

        # ── State ─────────────────────────────────────────────────────────────
        raw_state = gr.State("")

        # ── Wiring ───────────────────────────────────────────────────────────
        def run_generation(code, file_obj, dtype, ctx):

            md_acc = ""
            status = ""
            ppt_file = None

            for result in generate_documentation(
                code,
                file_obj,
                dtype,
                ctx
            ):

                # Safely handle variable-length outputs
                if isinstance(result, tuple):

                    if len(result) == 5:
                        md, raw, status, raw_state_val, ppt_file = result

                    elif len(result) == 4:
                        md, raw, status, raw_state_val = result
                        ppt_file = None

                    elif len(result) == 3:
                        md, raw, status = result
                        raw_state_val = raw
                        ppt_file = None

                    else:
                        continue

                else:
                    continue

                md_acc = md

                yield (
                    md,
                    raw,
                    status,
                    raw_state_val,
                    ppt_file
                )

            yield (
                md_acc,
                md_acc,
                status,
                md_acc,
                ppt_file
            )
        

        generate_btn.click(
            fn=run_generation,
            inputs=[code_input, file_upload, doc_type, extra_context],
            outputs=[md_out, raw_out, status_out, raw_state, ppt_download],
        )

        def do_download(content, dtype):
            if not content.strip():
                return None
            return create_download_file(content, dtype)

        download_btn.click(
            fn=do_download,
            inputs=[raw_state, doc_type],
            outputs=[download_file],
        )

        def do_clear():
            return "", None, "", "", "", "",None

        clear_btn.click(
            fn=do_clear,
            outputs=[code_input, file_upload, md_out, raw_out, status_out, raw_state,ppt_download],
        )

        # ── Sample loaders ────────────────────────────────────────────────────
        SAMPLE_PY = '''\
"""
auth.py — JWT Authentication module for a REST API.
"""
import jwt
import hashlib
import secrets
from datetime import datetime, timedelta
from typing import Optional

SECRET_KEY = "super-secret-key"
ALGORITHM  = "HS256"

def hash_password(password: str) -> str:
    """Hash a plaintext password using SHA-256 with a random salt."""
    salt = secrets.token_hex(16)
    hashed = hashlib.sha256((password + salt).encode()).hexdigest()
    return f"{salt}:{hashed}"

def verify_password(password: str, stored_hash: str) -> bool:
    """Verify a plaintext password against a stored hash.

    Args:
        password: The plaintext password to verify.
        stored_hash: The stored hash in 'salt:hash' format.

    Returns:
        True if the password matches, False otherwise.
    """
    salt, hashed = stored_hash.split(":")
    return hashlib.sha256((password + salt).encode()).hexdigest() == hashed

def create_access_token(user_id: int, expires_minutes: int = 30) -> str:
    """Create a signed JWT access token.

    Args:
        user_id: The user\'s unique identifier.
        expires_minutes: Token validity in minutes (default 30).

    Returns:
        A signed JWT string.

    Raises:
        ValueError: If user_id is not a positive integer.
    """
    if user_id <= 0:
        raise ValueError("user_id must be a positive integer")
    payload = {
        "sub": user_id,
        "exp": datetime.utcnow() + timedelta(minutes=expires_minutes),
        "iat": datetime.utcnow(),
    }
    return jwt.encode(payload, SECRET_KEY, algorithm=ALGORITHM)

def decode_token(token: str) -> Optional[dict]:
    """Decode and validate a JWT token.

    Args:
        token: The JWT string.

    Returns:
        The decoded payload dict, or None if invalid/expired.
    """
    try:
        return jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
    except jwt.PyJWTError:
        return None
'''

        SAMPLE_JS = '''\
/**
 * EventEmitter — a lightweight pub/sub event bus.
 * @module EventEmitter
 */
class EventEmitter {
  constructor() {
    /** @private */
    this._events = {};
  }

  /**
   * Register a listener for an event.
   * @param {string} event - Event name.
   * @param {Function} listener - Callback function.
   * @returns {EventEmitter} this (for chaining).
   */
  on(event, listener) {
    if (!this._events[event]) this._events[event] = [];
    this._events[event].push(listener);
    return this;
  }

  /**
   * Register a one-time listener.
   * @param {string} event - Event name.
   * @param {Function} listener - Callback, auto-removed after first call.
   * @returns {EventEmitter} this
   */
  once(event, listener) {
    const wrapper = (...args) => {
      listener(...args);
      this.off(event, wrapper);
    };
    return this.on(event, wrapper);
  }

  /**
   * Remove a specific listener.
   * @param {string} event
   * @param {Function} listener
   * @returns {EventEmitter} this
   */
  off(event, listener) {
    if (!this._events[event]) return this;
    this._events[event] = this._events[event].filter(l => l !== listener);
    return this;
  }

  /**
   * Emit an event, calling all registered listeners.
   * @param {string} event
   * @param {...*} args - Arguments passed to each listener.
   * @returns {boolean} True if listeners were called.
   */
  emit(event, ...args) {
    if (!this._events[event]?.length) return false;
    this._events[event].forEach(l => l(...args));
    return true;
  }
}

module.exports = EventEmitter;
'''

        SAMPLE_GO = '''\
// Package cache provides an in-memory LRU cache with TTL support.
package cache

import (
    "sync"
    "time"
)

// Entry holds a cached value and its expiration time.
type Entry struct {
    Value     interface{}
    ExpiresAt time.Time
}

// Cache is a thread-safe LRU cache.
type Cache struct {
    mu       sync.RWMutex
    items    map[string]*Entry
    maxItems int
    ttl      time.Duration
}

// New creates a new Cache with the given capacity and default TTL.
//
// Parameters:
//   maxItems - maximum number of items before eviction (must be > 0)
//   ttl      - default time-to-live for entries (0 = no expiration)
//
// Returns a pointer to the new Cache.
func New(maxItems int, ttl time.Duration) *Cache {
    return &Cache{items: make(map[string]*Entry), maxItems: maxItems, ttl: ttl}
}

// Set stores a key-value pair in the cache.
// If the cache is full the oldest entry is evicted.
func (c *Cache) Set(key string, value interface{}) {
    c.mu.Lock()
    defer c.mu.Unlock()
    exp := time.Time{}
    if c.ttl > 0 {
        exp = time.Now().Add(c.ttl)
    }
    c.items[key] = &Entry{Value: value, ExpiresAt: exp}
}

// Get retrieves a value by key.
// Returns (value, true) if found and not expired, else (nil, false).
func (c *Cache) Get(key string) (interface{}, bool) {
    c.mu.RLock()
    defer c.mu.RUnlock()
    e, ok := c.items[key]
    if !ok {
        return nil, false
    }
    if !e.ExpiresAt.IsZero() && time.Now().After(e.ExpiresAt) {
        return nil, false
    }
    return e.Value, true
}

// Delete removes a key from the cache.
func (c *Cache) Delete(key string) {
    c.mu.Lock()
    defer c.mu.Unlock()
    delete(c.items, key)
}

// Len returns the current number of items in the cache.
func (c *Cache) Len() int {
    c.mu.RLock()
    defer c.mu.RUnlock()
    return len(c.items)
}
'''

        sample_py.click(fn=lambda: SAMPLE_PY, outputs=[code_input])
        sample_js.click(fn=lambda: SAMPLE_JS, outputs=[code_input])
        sample_go.click(fn=lambda: SAMPLE_GO, outputs=[code_input])

    return demo


# ── Entry point ───────────────────────────────────────────────────────────────
if __name__ == "__main__":
    demo = build_ui()
    demo.launch(
        server_name="127.0.0.1",
        server_port=7860,
        share=False,
        show_error=True,
    )
